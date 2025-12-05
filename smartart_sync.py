import os
from itertools import zip_longest
from pathlib import Path

from pptx import Presentation
from pptx.diagram import SmartArt


def sync_smartart_text(smartart: SmartArt, text_list: list[str]) -> None:
    current_nodes = smartart.editable_nodes

    for node, text in zip_longest(current_nodes, text_list, fillvalue=None):
        if node is not None and text is not None:
            node.text = text
            continue

        if node is None and text is not None:
            smartart.add_node(text)


def sync_smartart_images(smartart: SmartArt, image_list: list[Path]) -> None:
    current_nodes = smartart.editable_nodes

    for node, image_path in zip_longest(current_nodes, image_list, fillvalue=None):
        if node is None:
            node = smartart.add_node()
            if not node.has_image_placeholder:
                smartart.remove_node(node)
                continue

        if not node.has_image_placeholder:
            continue

        if image_path is not None:
            node.image_path = str(image_path)


def remove_empty_smartart_nodes(smartart: SmartArt) -> None:
    current_nodes = smartart.editable_nodes

    for node in current_nodes:
        has_image = node.has_image_placeholder and node.image_path is not None
        has_text = bool(node.text.strip())

        if has_image or has_text:
            continue

        smartart.remove_node(node)


def main():
    source = Path(r"D:\Private\python_projects\pptx-slides-gpt-app\templates\single.pptx")
    presentation = Presentation(source)
    images = list((Path.cwd() / "images").glob("*.png"))
    text = [f"Item {i}" for i in range(len(images))]

    for slide in presentation.slides:
        for i, shape in enumerate(slide.shapes):
            print(f"Shape {i}")
            if not shape.has_smartart:
                continue

            smartart = shape.smartart
            sync_smartart_images(smartart, images)
            sync_smartart_text(smartart, text)
            remove_empty_smartart_nodes(smartart)

    output = source.parent / f"modified_{source.name}"
    presentation.save(output)
    os.startfile(output)


if __name__ == "__main__":
    main()
